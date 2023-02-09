import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import pytesseract
from sys import path
from os import listdir,path

# exec_path=os.getcwd()
# print(exec_path)
pytesseract.pytesseract.tesseract_cmd = r'C:\\Program Files\\Tesseract-OCR\\tesseract'
image_path = r"C:\Users\muhammad.ziada\Downloads\wetransfer_untitled-transfer_2022-11-05_0638"


def OCR(image_file):
    image = Image.open(f'{image_path}\\{image_file}')
    text = pytesseract.image_to_string(image).replace('|', 'I')
    return text.split('\n\n')

def Export_to_Slide(slides_paragraphs):
    left = top = Inches(1)
    height = Inches(6)
    width = Inches(4)

    font_name = 'Times New Roman'
    font_size = Pt(14)

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for slide in slides_paragraphs:
        try:
            paragraphs1=slide[0]
            paragraphs2=slide[1]

            slide = prs.slides.add_slide(blank_slide_layout)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            txBox2 = slide.shapes.add_textbox(
                left + width + Inches(.5), top, width, height)
            tf2 = txBox2.text_frame
            tf2.word_wrap = True

            for paragraph in paragraphs1:
                p = tf.add_paragraph()
                p.text = paragraph
                # p.font.bold = True
                p.font.name = font_name
                p.font.size = font_size
            
            for paragraph in paragraphs2:
                p2 = tf2.add_paragraph()
                p2.text = paragraph
                # p.font.bold = True
                p2.font.name = font_name
                p2.font.size = font_size
        except Exception:
            print(Exception)
            continue

    prs.save('test.pptx')


def OCR_1Paragraph(image_file):
    image = Image.open(f'{image_path}\\{image_file}')
    return pytesseract.image_to_string(image).replace('|', 'I')

def Export_to_Slide_1Column(slides):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[1]

    font_name = 'Times New Roman'
    font_size = Pt(14)

    for paragraph in slides:
        try:
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = ""
            content.text = paragraph
            content.font.name= font_name
            content.font.size= font_size
            content.word_wrap=True

        except Exception as e:
            print(e)
            continue

    prs.save('test2.pptx')


def main():
    files_list=listdir(image_path)
    # slides_paragraphs=[()]
    # for idx in range(0,len(files_list),2):
    #     file1=files_list[idx]
    #     file2=files_list[idx+1] if idx+1<len(files_list) else ""
    #     try:
    #         if file1.endswith('.jpeg'):
    #             print(file1)
    #             paragraphs1= OCR(file1)
    #         file2 = files_list[(idx + 1) % len(files_list)]
    #         if file2.endswith('.jpeg'):
    #             print(file2)
    #             paragraphs2= OCR(file2)
    #         paragraphs=(paragraphs1, paragraphs2)
    #         slides_paragraphs.append(paragraphs)
    #     except Exception:
    #         print(Exception)
    #         continue
    # Export_to_Slide(slides_paragraphs)

    slides_paragraphs = []
    for file in files_list:
        try:
            if file.endswith('.jpeg'):
                print(file)
                paragraph = OCR_1Paragraph(file)
                slides_paragraphs.append(paragraph)
        except Exception as e:
            print(e)
            continue
    Export_to_Slide_1Column(slides_paragraphs)


main()
